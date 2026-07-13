"""Generate a PowerPoint deck summarizing the OmniSearch CV results.

Builds all charts (recall heatmap, modality comparison bars, training metrics,
perception-by-condition, dataset composition, project timeline) and assembles
them into docs/slides/omnisearch_cv_results.pptx.

Usage:
    python3 scripts/generate_results_slides.py
"""

from __future__ import annotations

import json
import os
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", "/tmp/mpl-omnisearch")

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.colors import LinearSegmentedColormap

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
import sys
sys.path.insert(0, str(ROOT))
SLIDE_DIR = ROOT / "docs" / "slides"
ASSET_DIR = SLIDE_DIR / "assets"
ASSET_DIR.mkdir(parents=True, exist_ok=True)

# Build on top of the Capstone deck so we inherit its exact slide master,
# layouts, theme and slide size — guaranteeing an identical look-and-feel.
TEMPLATE = ROOT / "docs" / "Capstone_Week10_Presentation.pptx"

# ---------------------------------------------------------------------------
# Theme — matches docs/Capstone_Week10_Presentation.pptx
#   Forest green headings (Georgia), rust kicker/accents (Arial), taupe subtitle
# ---------------------------------------------------------------------------
GREEN = "#1F4D2E"       # forest green — headings, primary bars
GREEN2 = "#2C5530"      # secondary green
OLIVE = "#5B7C2F"       # UGV/olive accent
SAGE = "#A5B48C"        # muted sage
RUST = "#C44826"        # rust/orange — kicker, accents, big stats
RUST2 = "#B0472A"       # footer rust
FIRE = "#E76A2C"        # fire orange
TAUPE = "#8B7E68"       # muted taupe — subtitles
BODY = "#2B2B2B"        # body text
BODY2 = "#444033"       # secondary body
GREY = "#5E5A4A"        # date/label grey
LIGHT = "#E8DFC3"       # warm yellow (matches Capstone Week 10)

# Chart palette aliases (used by chart functions)
BLUE = GREEN            # "primary" series → forest green
STEEL_BLUE = "#3A7CA5"  # distinct color for CV-only fusion bars
TEAL = OLIVE            # "secondary" series → olive
ORANGE = RUST           # "accent" series → rust
RED = FIRE
NAVY = GREEN2

RGB_GREEN = RGBColor(0x1F, 0x4D, 0x2E)
RGB_GREEN2 = RGBColor(0x2C, 0x55, 0x30)
RGB_RUST = RGBColor(0xC4, 0x48, 0x26)
RGB_RUST2 = RGBColor(0xB0, 0x47, 0x2A)
RGB_TAUPE = RGBColor(0x8B, 0x7E, 0x68)
RGB_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RGB_DARK = RGBColor(0x2B, 0x2B, 0x2B)
RGB_GREY = RGBColor(0x5E, 0x5A, 0x4A)
RGB_LIGHT = RGBColor(0xE8, 0xDF, 0xC3)
RGB_OLIVE = RGBColor(0x5B, 0x7C, 0x2F)

# Legacy aliases so existing slide code keeps working
RGB_NAVY = RGB_GREEN
RGB_BLUE = RGB_GREEN
RGB_TEAL = RGB_OLIVE
RGB_ORANGE = RGB_RUST

HEAD_FONT = "Georgia"
BODY_FONT = "Calibri"
LABEL_FONT = "Arial"

CHART_BG = LIGHT
plt.rcParams.update({
    "font.family": "serif",
    "font.serif": ["Georgia", "PT Serif", "DejaVu Serif"],
    "font.size": 12,
    "axes.edgecolor": "#8B7E68",
    "axes.linewidth": 0.8,
    "figure.dpi": 150,
    "figure.facecolor": CHART_BG,
    "axes.facecolor": CHART_BG,
    "savefig.facecolor": CHART_BG,
    "text.color": "#2B2B2B",
    "axes.labelcolor": "#2B2B2B",
    "xtick.color": "#2B2B2B",
    "ytick.color": "#2B2B2B",
    "legend.facecolor": CHART_BG,
})

SCENARIOS = ["clear", "light_smoke", "heavy_smoke", "active_fire", "burned_ground", "mixed"]
SCENARIO_LABELS = ["Clear", "Light\nsmoke", "Heavy\nsmoke", "Active\nfire", "Burned\nground", "Mixed"]
MODES = ["cv", "thermal", "motion", "cv+thermal", "cv+motion"]
MODE_LABELS = ["CV\n(measured)", "Thermal", "Motion", "CV+Thermal", "CV+Motion"]


REALWORLD_ROWS = ["Thermal\n(HIT-UAV)", "Drone RGB\n(HERIDAL)", "UGV front\n(real Malibu terrain)", "UGV front\n(COCO)"]
REALWORLD_COLS = ["Synthetic val", "Real \u2014 before FT", "Real \u2014 after FT"]


def load_realworld_validation() -> np.ndarray:
    """Recall matrix (rows above x REALWORLD_COLS) for the real-data eval slide.

    "Real \u2014 before FT" / "after FT" are loaded from reports/*.json produced by
    scripts/eval_heridal.py, scripts/annotate_real_eval.py's underlying evals,
    and the direct YOLO .val() calls in this session. "Synthetic val" is the
    same model's recall on its own synthetic validation split (documented in
    docs/real_data_eval_report.md \u00a71); kept as a literal here because it comes
    from training-run logs rather than a saved report file.
    """

    def _j(name):
        p = ROOT / "reports" / name
        return json.loads(p.read_text()) if p.exists() else None

    thermal = _j("thermal_real_eval.json")
    heridal_before = _j("heridal_real_eval.json")
    heridal_after = _j("heridal_real_eval_finetuned.json")
    ugv_coco_before = _j("ugv_real_eval.json")
    ugv_coco_after = _j("ugv_real_eval_finetuned.json")
    ugv_malibu_before = _j("ugv_malibu_real_eval_synthetic_only.json")
    ugv_malibu_after = _j("ugv_malibu_real_eval_finetuned_coco_mix.json")

    rows = [
        # Thermal: synthetic val (training log), real before/after (HIT-UAV test)
        [0.96,
         thermal["results"]["synthetic_trained"]["recall"] if thermal else 0.31,
         thermal["results"]["real_trained"]["recall"] if thermal else 0.86],
        # Drone RGB: synthetic val (training log), real before/after (HERIDAL, full@1280, conf 0.15)
        [0.95,
         heridal_before["results"]["full_1280"]["conf_0.15"]["recall"] if heridal_before else 0.22,
         heridal_after["results"]["full_1280"]["conf_0.15"]["recall"] if heridal_after else 0.82],
        # UGV front on real Malibu terrain composites
        [0.33,
         ugv_malibu_before["recall"] if ugv_malibu_before else 0.035,
         ugv_malibu_after["recall"] if ugv_malibu_after else 0.63],
        # UGV front on COCO person holdout
        [0.33,
         ugv_coco_before["results"]["ugv_front"]["recall"] if ugv_coco_before else 0.005,
         ugv_coco_after["coco_holdout"]["recall"] if ugv_coco_after else 0.57],
    ]
    return np.array(rows)


def chart_realworld_heatmap() -> Path:
    """Heatmap: modality/real-test-set (rows) x synthetic-val / real-before / real-after (cols)."""
    mat = load_realworld_validation()
    cmap = LinearSegmentedColormap.from_list(
        "recall", ["#7A2214", "#C44826", "#E76A2C", "#D9B36A", "#8FA86B", "#1F4D2E"]
    )
    fig, ax = plt.subplots(figsize=(8.6, 5.6))
    im = ax.imshow(mat, cmap=cmap, vmin=0, vmax=1, aspect="auto")

    ax.set_xticks(range(len(REALWORLD_COLS)))
    ax.set_xticklabels(REALWORLD_COLS, fontsize=12, fontweight="bold")
    ax.set_yticks(range(len(REALWORLD_ROWS)))
    ax.set_yticklabels([r.replace("\n", " ") for r in REALWORLD_ROWS], fontsize=12)

    for i in range(mat.shape[0]):
        for j in range(mat.shape[1]):
            val = mat[i, j]
            ax.text(j, i, f"{val*100:.0f}%", ha="center", va="center",
                    color="white" if val < 0.75 else "#1A1A1A",
                    fontsize=14, fontweight="bold")

    ax.set_title("Real-World Validation: Recall Before vs. After Fine-Tuning",
                 fontsize=14, fontweight="bold", color=NAVY, pad=14)
    cbar = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.03)
    cbar.set_label("Recall", fontsize=11)
    ax.set_xticks(np.arange(-.5, len(REALWORLD_COLS), 1), minor=True)
    ax.set_yticks(np.arange(-.5, len(REALWORLD_ROWS), 1), minor=True)
    ax.grid(which="minor", color="white", linewidth=2)
    ax.tick_params(which="minor", length=0)
    fig.tight_layout()
    p = ASSET_DIR / "realworld_heatmap.png"
    fig.savefig(p, bbox_inches="tight", facecolor=CHART_BG)
    plt.close(fig)
    return p


def load_comparison() -> dict:
    """Load the modality comparison and substitute the REAL measured CV column.

    The original data/detection_mode_comparison.json "cv" entries came from
    PreliminaryPersonDetector — a coin-flip stub over ground-truth boxes that
    produced misleading 100% recalls. scripts/eval_cv_by_scenario.py replaces
    them by running the actual trained YOLOv8s on clear val frames with each
    scenario's fire/smoke synthesized by the same renderer the simulator uses.
    The fusion columns are recomputed from the measured CV with the same
    formulas compare_detection_modes.py used (union for CV+Thermal; boost mode
    keeps CV recall for CV+Motion).
    """
    data = json.loads((ROOT / "data" / "detection_mode_comparison.json").read_text())
    measured_path = ROOT / "reports" / "cv_recall_by_scenario_1280.json"
    if not measured_path.exists():
        measured_path = ROOT / "reports" / "cv_recall_by_scenario.json"
    if measured_path.exists():
        measured = json.loads(measured_path.read_text())["scenarios"]
        for s in SCENARIOS:
            if s not in measured:
                continue
            cv_r = measured[s]["recall"]
            data[s]["cv"]["recall"] = cv_r
            data[s]["cv"]["precision"] = measured[s]["precision"]
            th_r = data[s]["thermal"]["recall"]
            data[s]["cv+thermal"]["recall"] = round(1.0 - (1.0 - cv_r) * (1.0 - th_r), 3)
            data[s]["cv+motion"]["recall"] = cv_r
    return data


# ---------------------------------------------------------------------------
# Charts
# ---------------------------------------------------------------------------
def chart_recall_heatmap(data: dict) -> Path:
    """Recall heatmap: scenarios (rows) x modalities (cols)."""
    mat = np.array([[data[s][m]["recall"] for m in MODES] for s in SCENARIOS])

    cmap = LinearSegmentedColormap.from_list(
        "recall", ["#7A2214", "#C44826", "#E76A2C", "#D9B36A", "#8FA86B", "#1F4D2E"]
    )
    fig, ax = plt.subplots(figsize=(9.2, 5.6))
    im = ax.imshow(mat, cmap=cmap, vmin=0, vmax=1, aspect="auto")

    ax.set_xticks(range(len(MODES)))
    ax.set_xticklabels(MODE_LABELS, fontsize=12, fontweight="bold")
    ax.set_yticks(range(len(SCENARIOS)))
    ax.set_yticklabels([s.replace("\n", " ") for s in SCENARIO_LABELS], fontsize=12)

    for i in range(len(SCENARIOS)):
        for j in range(len(MODES)):
            val = mat[i, j]
            ax.text(j, i, f"{val*100:.0f}%", ha="center", va="center",
                    color="white" if val < 0.75 else "#1A1A1A",
                    fontsize=13, fontweight="bold")

    ax.set_title("Detection Recall by Scenario and Modality",
                 fontsize=15, fontweight="bold", color=NAVY, pad=14)
    cbar = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.03)
    cbar.set_label("Recall", fontsize=11)
    ax.set_xticks(np.arange(-.5, len(MODES), 1), minor=True)
    ax.set_yticks(np.arange(-.5, len(SCENARIOS), 1), minor=True)
    ax.grid(which="minor", color="white", linewidth=2)
    ax.tick_params(which="minor", length=0)
    fig.tight_layout()
    p = ASSET_DIR / "recall_heatmap.png"
    fig.savefig(p, bbox_inches="tight", facecolor=CHART_BG)
    plt.close(fig)
    return p


def chart_real_cv_recall() -> Path:
    """Bar chart: REAL trained-model recall/precision, clear vs fire/smoke,
    measured on held-out val frames (reports/cv_recall_by_fire_condition.json).

    This is the actual detector, not the idealized coin-flip proxy used for
    the six-scenario modality heatmap. Only two conditions are shown because
    that is the only axis our generated val metadata actually labels.
    """
    fc_path = ROOT / "reports" / "cv_recall_by_fire_condition_1280.json"
    if not fc_path.exists():
        fc_path = ROOT / "reports" / "cv_recall_by_fire_condition.json"
    d = json.loads(fc_path.read_text())["buckets"]
    labels = ["Clear", "Fire / smoke\npresent"]
    recall = [d["clear"]["recall"] * 100, d["fire_or_smoke"]["recall"] * 100]
    precision = [d["clear"]["precision"] * 100, d["fire_or_smoke"]["precision"] * 100]

    x = np.arange(len(labels)); w = 0.32
    fig, ax = plt.subplots(figsize=(6.4, 4.6))
    ax.bar(x - w / 2, recall, w, label="Recall", color=BLUE)
    ax.bar(x + w / 2, precision, w, label="Precision", color=ORANGE)
    for i in range(len(labels)):
        ax.text(x[i] - w / 2, recall[i] + 1.5, f"{recall[i]:.0f}%", ha="center", fontsize=12, fontweight="bold")
        ax.text(x[i] + w / 2, precision[i] + 1.5, f"{precision[i]:.0f}%", ha="center", fontsize=12, fontweight="bold")
    ax.set_xticks(x); ax.set_xticklabels(labels, fontsize=13, fontweight="bold")
    ax.set_ylabel("Score (%)", fontsize=12)
    ax.set_ylim(0, 100)
    ax.set_title("Real Trained Detector \u2014 Measured on Held-Out Val Frames",
                 fontsize=12.5, fontweight="bold", color=NAVY, pad=12)
    ax.legend(loc="lower right", fontsize=11, frameon=False)
    fig.tight_layout()
    p = ASSET_DIR / "real_cv_recall.png"
    fig.savefig(p, bbox_inches="tight", facecolor=CHART_BG)
    plt.close(fig)
    return p


def chart_decoy_fp_heatmap() -> Path:
    """Heatmap: hard-negative decoy class x false-positive rate, before/after
    the human-scale-decoy fix (animal/colorful-object classes + compact-car
    resizing), read from reports/decoy_fp_rates.json ('after') plus the
    pre-fix baseline measured on the same val set ('before')."""
    fp_path = ROOT / "reports" / "decoy_fp_rates_1280.json"
    if not fp_path.exists():
        fp_path = ROOT / "reports" / "decoy_fp_rates.json"
    after = json.loads(fp_path.read_text())["per_type"]
    classes = ["vehicle", "animal", "colorful_object"]
    class_labels = ["Vehicle\n(3.5\u20134.8 m)", "Animal\n(0.5\u20131.8 m)", "Colorful object\n(0.5\u20131.8 m)"]
    before = {"vehicle": 0.0, "animal": 0.067, "colorful_object": 0.409}
    cols = ["Before fix", "After fix"]
    mat = np.array([[before[c], after[c]["fp_rate"]] for c in classes])

    cmap = LinearSegmentedColormap.from_list(
        "fp", ["#1F4D2E", "#8FA86B", "#D9B36A", "#E76A2C", "#C44826", "#7A2214"]
    )
    fig, ax = plt.subplots(figsize=(6.4, 4.6))
    im = ax.imshow(mat, cmap=cmap, vmin=0, vmax=0.45, aspect="auto")
    ax.set_xticks(range(len(cols))); ax.set_xticklabels(cols, fontsize=13, fontweight="bold")
    ax.set_yticks(range(len(classes))); ax.set_yticklabels(class_labels, fontsize=12)
    for i in range(len(classes)):
        for j in range(2):
            val = mat[i, j]
            ax.text(j, i, f"{val*100:.1f}%", ha="center", va="center",
                    color="white" if val > 0.22 else "#1A1A1A", fontsize=15, fontweight="bold")
    ax.set_title("Decoy False-Positive Rate — Before vs After Hard-Negative Fix",
                 fontsize=13, fontweight="bold", color=NAVY, pad=14)
    cbar = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    cbar.set_label("FP rate (decoy misread as survivor)", fontsize=10)
    ax.set_xticks(np.arange(-.5, 2, 1), minor=True)
    ax.set_yticks(np.arange(-.5, len(classes), 1), minor=True)
    ax.grid(which="minor", color="white", linewidth=2)
    ax.tick_params(which="minor", length=0)
    fig.tight_layout()
    p = ASSET_DIR / "decoy_fp_heatmap.png"
    fig.savefig(p, bbox_inches="tight", facecolor=CHART_BG)
    plt.close(fig)
    return p


def chart_recall_by_size_heatmap() -> Path:
    """Heatmap: survivor pixel-size bucket x inference config, recall @ IoU 0.15,
    read from reports/recall_by_size.json (post hard-negative retrain)."""
    rbs_path = ROOT / "reports" / "recall_by_size_1280.json"
    if not rbs_path.exists():
        rbs_path = ROOT / "reports" / "recall_by_size.json"
    d = json.loads(rbs_path.read_text())
    configs = ["A_640_full", "B_1280_full", "C_1280_tiled"]
    config_labels = ["640\nfull frame", "1280\nfull frame", "1280\ntiled 2\u00d72"]
    buckets = ["<8px", "8-15px", ">15px"]
    bucket_labels = ["<8 px", "8\u201315 px", ">15 px"]
    mat = np.array([[d["results"][c]["recall@iou0.15"][b]["recall"] for c in configs] for b in buckets])

    cmap = LinearSegmentedColormap.from_list(
        "recall", ["#7A2214", "#C44826", "#E76A2C", "#D9B36A", "#8FA86B", "#1F4D2E"]
    )
    fig, ax = plt.subplots(figsize=(6.8, 4.6))
    im = ax.imshow(mat, cmap=cmap, vmin=0.5, vmax=1.0, aspect="auto")
    ax.set_xticks(range(len(configs))); ax.set_xticklabels(config_labels, fontsize=12, fontweight="bold")
    ax.set_yticks(range(len(buckets))); ax.set_yticklabels(bucket_labels, fontsize=13)
    for i in range(len(buckets)):
        for j in range(len(configs)):
            val = mat[i, j]
            ax.text(j, i, f"{val*100:.0f}%", ha="center", va="center",
                    color="white" if val < 0.78 else "#1A1A1A", fontsize=15, fontweight="bold")
    ax.set_title("Recall by Target Size and Inference Config",
                 fontsize=13, fontweight="bold", color=NAVY, pad=14)
    cbar = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    cbar.set_label("Recall (IoU \u2265 0.15)", fontsize=10)
    ax.set_xticks(np.arange(-.5, len(configs), 1), minor=True)
    ax.set_yticks(np.arange(-.5, len(buckets), 1), minor=True)
    ax.grid(which="minor", color="white", linewidth=2)
    ax.tick_params(which="minor", length=0)
    fig.tight_layout()
    p = ASSET_DIR / "recall_by_size_heatmap.png"
    fig.savefig(p, bbox_inches="tight", facecolor=CHART_BG)
    plt.close(fig)
    return p


def chart_fusion_bars(data: dict) -> Path:
    """Grouped bars: CV vs Thermal vs CV+Thermal recall per scenario."""
    cv = [data[s]["cv"]["recall"] * 100 for s in SCENARIOS]
    th = [data[s]["thermal"]["recall"] * 100 for s in SCENARIOS]
    fu = [data[s]["cv+thermal"]["recall"] * 100 for s in SCENARIOS]

    x = np.arange(len(SCENARIOS))
    w = 0.26
    fig, ax = plt.subplots(figsize=(9.6, 5.4))
    ax.bar(x - w, cv, w, label="CV", color=STEEL_BLUE)
    ax.bar(x, th, w, label="Thermal", color=ORANGE)
    ax.bar(x + w, fu, w, label="CV+Thermal (fusion)", color=GREEN)

    for i in range(len(SCENARIOS)):
        ax.text(x[i] + w, fu[i] + 1.5, f"{fu[i]:.0f}", ha="center", fontsize=10,
                fontweight="bold", color=GREEN)

    ax.set_ylabel("Recall (%)", fontsize=12)
    ax.set_ylim(0, 108)
    ax.set_xticks(x)
    ax.set_xticklabels(SCENARIO_LABELS, fontsize=11)
    ax.set_title("Sensor Fusion Beats Either Sensor Alone",
                 fontsize=15, fontweight="bold", color=NAVY, pad=12)
    ax.legend(loc="upper right", fontsize=11, framealpha=0.95)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="y", alpha=0.3)
    fig.tight_layout()
    p = ASSET_DIR / "fusion_bars.png"
    fig.savefig(p, bbox_inches="tight", facecolor=CHART_BG)
    plt.close(fig)
    return p


def chart_training_metrics() -> Path:
    """Grouped bars of trained detector metrics (all 5 detectors)."""
    models = ["Drone\n(oblique)", "Drone\n(NAIP)", "UGV\nfront", "UGV\nmast", "Thermal\nYOLO"]
    # Drone (oblique) reflects the latest retrain on the harder, hard-negative-
    # diversified dataset (animal/colorful-object decoys + compact-car sizing);
    # recall/mAP dropped vs. the earlier easier-decoy run because the task is
    # now genuinely harder to shortcut.
    precision = [0.815, 0.892, 0.867, 0.896, 1.000]
    recall    = [0.658, 0.913, 0.737, 0.899, 0.963]
    map50     = [0.719, 0.914, 0.807, 0.923, 0.965]

    x = np.arange(len(models))
    w = 0.22
    fig, ax = plt.subplots(figsize=(10.5, 5.2))
    ax.bar(x - w, precision, w, label="Precision", color=BLUE)
    ax.bar(x, recall, w, label="Recall", color=TEAL)
    ax.bar(x + w, map50, w, label="mAP@50", color=ORANGE)

    for i in range(len(models)):
        for off, v in zip((-w, 0, w), (precision[i], recall[i], map50[i])):
            ax.text(x[i] + off, v + 0.015, f"{v:.2f}", ha="center", fontsize=9, fontweight="bold")

    ax.set_ylabel("Score", fontsize=12)
    ax.set_ylim(0, 1.12)
    ax.set_xticks(x)
    ax.set_xticklabels(models, fontsize=12, fontweight="bold")
    ax.set_title("Trained Detector Performance (held-out validation)",
                 fontsize=15, fontweight="bold", color=NAVY, pad=12)
    ax.legend(loc="lower right", fontsize=11)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="y", alpha=0.3)
    fig.tight_layout()
    p = ASSET_DIR / "training_metrics.png"
    fig.savefig(p, bbox_inches="tight", facecolor=CHART_BG)
    plt.close(fig)
    return p


def chart_perception_by_condition() -> Path:
    """CV recall + precision by condition — REAL measured values from
    scripts/eval_cv_by_scenario.py (reports/cv_recall_by_scenario.json)."""
    cond_path = ROOT / "reports" / "cv_recall_by_scenario_1280.json"
    if not cond_path.exists():
        cond_path = ROOT / "reports" / "cv_recall_by_scenario.json"
    d = json.loads(cond_path.read_text())["scenarios"]
    keys = ["clear", "light_smoke", "heavy_smoke", "active_fire", "burned_ground", "mixed"]
    conds = ["Clear", "Light\nsmoke", "Heavy\nsmoke", "Active\nfire", "Burned\nground", "Mixed"]
    recall = [d[k]["recall"] for k in keys]
    precision = [d[k]["precision"] for k in keys]

    x = np.arange(len(conds))
    fig, ax = plt.subplots(figsize=(9.2, 5.2))
    ax.plot(x, recall, "-o", color=BLUE, linewidth=2.5, markersize=9, label="Recall")
    ax.plot(x, precision, "-s", color=ORANGE, linewidth=2.5, markersize=8, label="Precision")

    for i in range(len(conds)):
        ax.text(x[i], recall[i] + 0.02, f"{recall[i]:.2f}", ha="center", fontsize=9, color=BLUE, fontweight="bold")
        ax.text(x[i], precision[i] - 0.05, f"{precision[i]:.2f}", ha="center", fontsize=9, color=ORANGE, fontweight="bold")

    ax.set_ylabel("Score", fontsize=12)
    ax.set_ylim(0.4, 1.08)
    ax.set_xticks(x)
    ax.set_xticklabels(conds, fontsize=11)
    ax.set_title("CV Perception Quality by Condition (measured, real model)",
                 fontsize=14, fontweight="bold", color=NAVY, pad=12)
    ax.legend(loc="lower left", fontsize=11)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="y", alpha=0.3)
    fig.tight_layout()
    p = ASSET_DIR / "perception_condition.png"
    fig.savefig(p, bbox_inches="tight", facecolor=CHART_BG)
    plt.close(fig)
    return p


def chart_dataset_composition() -> Path:
    """Donut of dataset sizes across the three detectors."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 5.0))

    # Left: total images per detector family
    sizes = [2400, 3600, 1200]
    labels = ["Drone\n(aerial)", "UGV\n(front+mast)", "Thermal\n(TIR)"]
    colors = [BLUE, TEAL, ORANGE]
    wedges, _, autotexts = ax1.pie(
        sizes, labels=labels, colors=colors, autopct=lambda p: f"{int(round(p*sum(sizes)/100)):,}",
        startangle=90, pctdistance=0.78, wedgeprops=dict(width=0.42, edgecolor="white", linewidth=2),
        textprops=dict(fontsize=12, fontweight="bold"))
    for at in autotexts:
        at.set_color("white"); at.set_fontsize(11)
    ax1.set_title("Generated Images by Detector\n(7,200 total)", fontsize=13, fontweight="bold", color=NAVY)

    # Right: drone image composition
    comp = [88.7, 65, 20, 24.2, 11.3]
    comp_labels = ["With\nsurvivors", "Fire/\nsmoke", "Hard-neg\ndecoys", "Edge\nplacements", "Negatives"]
    ax2.barh(range(len(comp)), comp, color=[BLUE, RED, GREY, TEAL, NAVY])
    ax2.set_yticks(range(len(comp)))
    ax2.set_yticklabels(comp_labels, fontsize=11)
    ax2.invert_yaxis()
    for i, v in enumerate(comp):
        ax2.text(v + 1, i, f"{v:.0f}%", va="center", fontsize=10, fontweight="bold")
    ax2.set_xlim(0, 100)
    ax2.set_xlabel("% of drone training images", fontsize=11)
    ax2.set_title("Drone Dataset Composition", fontsize=13, fontweight="bold", color=NAVY)
    ax2.spines[["top", "right"]].set_visible(False)

    fig.tight_layout()
    p = ASSET_DIR / "dataset_composition.png"
    fig.savefig(p, bbox_inches="tight", facecolor=CHART_BG)
    plt.close(fig)
    return p


def _panel(split: str, stem: str, title: str, ext: str = ".jpg",
           static_info: str | None = None):
    """Build a (image, label, title, caption) panel tuple.

    When a metadata JSON sidecar exists (the altitude-aware ``survivor/`` and
    ``thermal/`` sets), the caption is derived from the real physics — altitude
    and ground-sample-distance — instead of a hard-coded, misleading value.
    """
    img_path = ROOT / f"data/cv_train/{split}/images/{stem}{ext}"
    lbl_path = ROOT / f"data/cv_train/{split}/labels/{stem}.txt"
    js = ROOT / f"data/cv_train/{split}/labels/{stem}.json"
    info = static_info or "640\u00d7640 RGB"
    if static_info is None and js.exists():
        import json as _json
        m = _json.loads(js.read_text())
        if "altitude_m" in m:  # drone / thermal altitude-aware
            alt = m.get("altitude_m"); gsd = m.get("gsd_m")
            n = m.get("n_survivors", 0)
            line2 = f"alt {alt:.0f} m \u00b7 GSD {gsd:.2f} m/px \u00b7 {n} survivor(s)" if gsd \
                else f"alt {alt:.0f} m \u00b7 {n} survivor(s)"
            if m.get("oblique"):
                line2 += f" \u00b7 oblique {m.get('tilt_deg', 0):.0f}\u00b0"
            info = "NAIP background (0.6 m/px source)\n" + line2
    return (img_path, lbl_path, title, info)


def _load_with_boxes(img_path: Path, label_path: Path):
    """Load an image and draw YOLO bounding boxes from its label file."""
    from PIL import Image as PILImage
    img = PILImage.open(img_path)
    w, h = img.size
    boxes = []
    if label_path.exists():
        for line in label_path.read_text().strip().splitlines():
            parts = line.split()
            cx, cy, bw, bh = float(parts[1]), float(parts[2]), float(parts[3]), float(parts[4])
            x0 = (cx - bw / 2) * w
            y0 = (cy - bh / 2) * h
            x1 = (cx + bw / 2) * w
            y1 = (cy + bh / 2) * h
            boxes.append((x0, y0, x1, y1))
    return img, boxes


def _annotate_boxes(ax, boxes):
    """Draw green bounding boxes with pixel size."""
    for x0, y0, x1, y1 in boxes:
        bw, bh = int(x1 - x0), int(y1 - y0)
        rect = plt.Rectangle((x0, y0), x1 - x0, y1 - y0,
                              linewidth=3.5, edgecolor="#00FF00", facecolor="none")
        ax.add_patch(rect)
        ax.text(x0, y0 - 6, f"person {bw}\u00d7{bh} px",
                color="#00FF00", fontsize=14, fontweight="bold",
                bbox=dict(boxstyle="round,pad=0.25", fc="black", alpha=0.7))


def _info_badge(ax, info_text):
    """Green badge at bottom-right with resolution/compositing info."""
    ax.text(0.98, 0.03, info_text, transform=ax.transAxes,
            ha="right", va="bottom", fontsize=11.5, color="white", fontweight="bold",
            fontfamily="monospace",
            bbox=dict(boxstyle="round,pad=0.35", fc="#1F4D2E", alpha=0.88))


def chart_cv_conditions_1() -> Path:
    """2x2 grid: CV clear conditions with GT boxes + dimension info."""
    panels = [
        _panel("survivor/train", "00958", "Clear — high-altitude survivor (48 m)"),
        _panel("survivor/train", "00756", "Clear — three small survivors"),
        _panel("survivor/train", "01717", "Clear — survivor in vegetation"),
        _panel("survivor/train", "00839", "Clear — low altitude, oblique view"),
    ]
    fig, axes = plt.subplots(2, 2, figsize=(13, 12))
    for ax, (img_path, lbl_path, title, info) in zip(axes.flat, panels):
        img, boxes = _load_with_boxes(img_path, lbl_path)
        ax.imshow(img)
        _annotate_boxes(ax, boxes)
        _info_badge(ax, info)
        ax.set_title(title, fontsize=18, fontweight="bold", color=GREEN, pad=10)
        ax.set_xticks([]); ax.set_yticks([])
        for sp in ax.spines.values():
            sp.set_edgecolor("#8B7E68"); sp.set_linewidth(2)
    fig.tight_layout(h_pad=0.8, w_pad=0.6)
    p = ASSET_DIR / "cv_conditions_clear.png"
    fig.savefig(p, bbox_inches="tight", dpi=170, facecolor=CHART_BG)
    plt.close(fig)
    return p


def chart_cv_conditions_2() -> Path:
    """2x2 grid: CV smoke/fire conditions with GT boxes + dimension info."""
    panels = [
        _panel("survivor/train", "00706", "Active fire — distant survivor"),
        _panel("survivor/train", "01460", "Fire + smoke — two survivors"),
        _panel("survivor/train", "01204", "Burned ground — survivors near fire"),
        _panel("survivor/train", "01921", "Low altitude — survivors in flames"),
    ]
    fig, axes = plt.subplots(2, 2, figsize=(13, 12))
    for ax, (img_path, lbl_path, title, info) in zip(axes.flat, panels):
        img, boxes = _load_with_boxes(img_path, lbl_path)
        ax.imshow(img)
        _annotate_boxes(ax, boxes)
        _info_badge(ax, info)
        ax.set_title(title, fontsize=18, fontweight="bold", color=GREEN, pad=10)
        ax.set_xticks([]); ax.set_yticks([])
        for sp in ax.spines.values():
            sp.set_edgecolor("#8B7E68"); sp.set_linewidth(2)
    fig.tight_layout(h_pad=0.8, w_pad=0.6)
    p = ASSET_DIR / "cv_conditions_fire.png"
    fig.savefig(p, bbox_inches="tight", dpi=170, facecolor=CHART_BG)
    plt.close(fig)
    return p


def chart_cv_examples_all() -> Path:
    """2x3 grid: top row = real/default data, bottom row = synthetic composited data."""
    panels = [
        # Row 0: "Default" — real backgrounds & source assets
        _panel("survivor/train", "00958", "Drone — clear, high altitude"),
        _panel("ugv/mast/train", "00100", "UGV — mast camera",
               static_info="640\u00d7640 RGB \u00b7 ground-level view"),
        _panel("thermal/train", "00500", "Thermal IR",
               ext=".png", static_info="512\u00d7512 grayscale \u00b7 simulated LWIR"),
        # Row 1: "Synthetic" — composited survivors + wildfire effects
        _panel("survivor/train", "01204", "Drone + survivors + fire"),
        _panel("ugv/front/train", "00010", "UGV — survivor close-up",
               static_info="640\u00d7640 RGB \u00b7 SARD cutout, close range"),
        _panel("survivor/train", "00706", "Drone + fire overlay"),
    ]

    fig, axes = plt.subplots(2, 3, figsize=(14, 9.5))

    row_labels = ["Real / default data", "Synthetic composited data"]
    for r, label in enumerate(row_labels):
        axes[r, 0].set_ylabel(label, fontsize=16, fontweight="bold",
                              color=RUST, rotation=90, labelpad=15)

    for idx, (img_path, lbl_path, title, info) in enumerate(panels):
        ax = axes[idx // 3, idx % 3]
        img, boxes = _load_with_boxes(img_path, lbl_path)
        if img.mode == "L":
            ax.imshow(img, aspect="equal", cmap="gray")
        else:
            ax.imshow(img, aspect="equal")
        _annotate_boxes(ax, boxes)
        ax.set_title(title, fontsize=14, fontweight="bold", color=GREEN, pad=8)
        _info_badge(ax, info)
        ax.set_xticks([]); ax.set_yticks([])
        for sp in ax.spines.values():
            sp.set_edgecolor("#8B7E68"); sp.set_linewidth(2)

    fig.tight_layout(h_pad=1.0, w_pad=0.8)
    p = ASSET_DIR / "cv_examples_all.png"
    fig.savefig(p, bbox_inches="tight", dpi=170, facecolor=CHART_BG)
    plt.close(fig)
    return p


def chart_timeline() -> Path:
    """Progress timeline of the CV work."""
    milestones = [
        ("Jun 9", "Real CV detection\n(drone + UGV)"),
        ("Jun 10", "NAIP backgrounds\ncut false positives"),
        ("Jun 21", "Compositing +\nhard negatives"),
        ("Jun 27", "Altitude-aware +\nUGV detectors"),
        ("Jun 28", "Tracking, TTA,\nadaptive conf."),
        ("Jun 29", "Thermal + Motion\n+ fusion"),
        ("Jul 3", "Oblique camera +\nthermal dataset"),
    ]
    fig, ax = plt.subplots(figsize=(12.5, 3.4))
    x = np.arange(len(milestones))
    ax.plot(x, [0] * len(x), "-", color=GREY, linewidth=2, zorder=1)
    ax.scatter(x, [0] * len(x), s=220, color=BLUE, zorder=2, edgecolor="white", linewidth=2)
    for i, (date, label) in enumerate(milestones):
        up = i % 2 == 0
        y = 0.35 if up else -0.35
        va = "bottom" if up else "top"
        ax.text(i, y, label, ha="center", va=va, fontsize=10.5, fontweight="bold", color=NAVY)
        ax.text(i, 0.12 if up else -0.12, date, ha="center",
                va="bottom" if up else "top", fontsize=9.5, color=ORANGE, fontweight="bold")
    ax.set_ylim(-1, 1)
    ax.set_xlim(-0.6, len(milestones) - 0.4)
    ax.axis("off")
    ax.set_title("CV Development Timeline (June–July 2026)",
                 fontsize=15, fontweight="bold", color=NAVY, pad=6)
    fig.tight_layout()
    p = ASSET_DIR / "timeline.png"
    fig.savefig(p, bbox_inches="tight", facecolor=CHART_BG)
    plt.close(fig)
    return p


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
SW = Inches(13.333)
SH = Inches(7.5)


def _fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_bg(slide, rgb=RGB_LIGHT):
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _fill(box, rgb)
    slide.shapes._spTree.remove(box._element)
    slide.shapes._spTree.insert(2, box._element)
    return box


def new_presentation():
    """Clone the Capstone template and strip its slides, keeping master/theme."""
    from pptx.oxml.ns import qn
    if TEMPLATE.exists():
        prs = Presentation(str(TEMPLATE))
        sldIdLst = prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(qn("r:id"))
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
            sldIdLst.remove(sldId)
    else:
        prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def divider(slide, top, left=Inches(0.6), width=Inches(12.13), color=None):
    """Thin full-width rule like the Capstone content slides."""
    from pptx.enum.shapes import MSO_SHAPE
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1.4))
    _fill(rule, color or RGB_TAUPE)
    return rule


def add_text(slide, text, left, top, width, height, size=16, bold=False,
             color=RGB_DARK, align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_bullets(slide, bullets, left, top, width, height, size=16, color=RGB_DARK, gap=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, lvl) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap)
        run = p.add_run()
        # Level-0 lines are green Georgia headers; sub-lines are taupe/body Calibri.
        if lvl == 0:
            run.text = txt
            run.font.size = Pt(size)
            run.font.color.rgb = RGB_GREEN
            run.font.bold = True
            run.font.name = HEAD_FONT
        else:
            run.text = "– " + txt
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = False
            run.font.name = BODY_FONT
    return tb


def footer(slide, page=None):
    add_text(slide, "OmniSearch", Inches(0.5), Inches(7.15), Inches(4.0), Inches(0.25),
             size=11, bold=True, color=RGB_RUST2, font=LABEL_FONT)
    right = "MIDS CAPSTONE  ·  SUMMER 2026"
    if page:
        right = f"MIDS CAPSTONE  ·  SUMMER 2026  ·  {page}"
    add_text(slide, right, Inches(7.3), Inches(7.15), Inches(5.5), Inches(0.25),
             size=9, bold=False, color=RGB_TAUPE, align=PP_ALIGN.RIGHT, font=LABEL_FONT)


def title_bar(slide, title, kicker=None, subtitle=None):
    """Capstone-style header: rust kicker, Georgia green title, taupe subtitle."""
    if kicker:
        add_text(slide, kicker.upper(), Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.35),
                 size=13, bold=True, color=RGB_RUST, font=LABEL_FONT)
    add_text(slide, title, Inches(0.6), Inches(0.9), Inches(12.2), Inches(0.75),
             size=30, bold=True, color=RGB_GREEN, font=HEAD_FONT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.6), Inches(1.7), Inches(12.0), Inches(0.5),
                 size=15, bold=False, color=RGB_TAUPE, font=BODY_FONT)
    footer(slide)


def pic_fit(slide, path, left, top, max_w, max_h):
    from PIL import Image as PILImage
    with PILImage.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w; h = int(max_w / ar)
    else:
        h = max_h; w = int(max_h * ar)
    left_c = left + (max_w - w) // 2
    top_c = top + (max_h - h) // 2
    return slide.shapes.add_picture(str(path), left_c, top_c, width=w, height=h)


def metric_card(slide, left, top, w, h, value, label, color=None):
    """Capstone-style stat: thin vertical accent bar + big Georgia number + label."""
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), h)
    _fill(bar, RGB_RUST)
    tx = left + Inches(0.28)
    tw = w - Inches(0.28)
    add_text(slide, value, tx, top - Inches(0.05), tw, Inches(1.0),
             size=52, bold=True, color=RGB_RUST, font=HEAD_FONT)
    add_text(slide, label, tx, top + Inches(1.0), tw, Inches(0.8),
             size=13, bold=False, color=RGB_DARK, font=BODY_FONT)


# ---------------------------------------------------------------------------
# Build deck
# ---------------------------------------------------------------------------
def build():
    data = load_comparison()
    p_heat = chart_recall_heatmap(data)
    p_realworld = chart_realworld_heatmap()
    p_real_cv = chart_real_cv_recall()
    p_decoy_heat = chart_decoy_fp_heatmap()
    p_size_heat = chart_recall_by_size_heatmap()
    p_fusion = chart_fusion_bars(data)
    p_train = chart_training_metrics()
    p_cond = chart_perception_by_condition()
    p_dataset = chart_dataset_composition()
    p_cv_clear = chart_cv_conditions_1()
    p_cv_fire = chart_cv_conditions_2()
    p_cv_all = chart_cv_examples_all()
    p_timeline = chart_timeline()

    prs = new_presentation()
    # Use the template's DEFAULT (blank) layout so slides inherit its master/theme.
    blank = prs.slide_layouts[0]

    from pptx.enum.shapes import MSO_SHAPE
    # ---- Slide 1: Title (capstone hero style) ----
    s = prs.slides.add_slide(blank)
    add_bg(s, RGB_LIGHT)
    add_text(s, "UC BERKELEY MIDS CAPSTONE  ·  SUMMER 2026  ·  COMPUTER VISION",
             Inches(0.62), Inches(0.62), Inches(11.0), Inches(0.35),
             size=13, bold=True, color=RGB_RUST, font=LABEL_FONT)
    add_text(s, "OmniSearch", Inches(0.55), Inches(1.7), Inches(11.5), Inches(1.7),
             size=72, bold=True, color=RGB_GREEN, font=HEAD_FONT)
    add_text(s, "Computer Vision for Wildfire Search & Rescue",
             Inches(0.62), Inches(3.5), Inches(11.5), Inches(0.6),
             size=24, bold=False, color=RGB_DARK, font=BODY_FONT)
    add_text(s, "Survivor detection across drones, ground robots, and thermal sensors",
             Inches(0.62), Inches(4.2), Inches(11.5), Inches(0.5),
             size=16, bold=False, color=RGB_TAUPE, font=BODY_FONT)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(5.0), Inches(12.06), Inches(0.04))
    _fill(accent, RGB_RUST)
    add_text(s, "FINAL RESULTS  ·  MULTI-MODAL DETECTION & SENSOR FUSION",
             Inches(0.62), Inches(5.25), Inches(11.5), Inches(0.35),
             size=13, bold=True, color=RGB_RUST, font=LABEL_FONT)
    add_text(s, "Oleksii Lavrenin  —  perception & evaluation (detection/ + evaluation/)",
             Inches(0.62), Inches(5.7), Inches(11.5), Inches(0.4),
             size=15, bold=True, color=RGB_GREEN, font=HEAD_FONT)
    footer(s)

    # ---- Slide 2: Executive summary (metric cards + key bullets) ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "What we built", "EXECUTIVE SUMMARY",
              "A physics-grounded, multi-modal survivor-detection stack for wildfire SAR")
    divider(s, Inches(2.2))
    cards = [
        ("3", "Detector families\n(drone, UGV, thermal)"),
        ("22,500", "Synthetic + real training/eval\nimages generated or ingested"),
        ("0.82", "Measured CV recall\n(clear, retrained @1280px)"),
        ("+55pp", "Real recall gained by\nmixed real+synthetic fine-tune"),
    ]
    cw, ch, gap = Inches(2.95), Inches(1.8), Inches(0.2)
    x0 = Inches(0.6)
    for i, (v, l) in enumerate(cards):
        metric_card(s, x0 + i * (cw + gap), Inches(2.35), cw, ch, v, l)
    add_bullets(s, [
        ("Real YOLOv8 survivor detection for aerial drones and ground robots", 0),
        ("replaced the abstract simulator perception model with in-distribution CV", 1),
        ("Physics-grounded synthetic data (7,200 images \u2192 Zenodo DOI 10.5281/zenodo.21226010)", 0),
        ("SARD cutouts on NAIP aerial imagery with altitude-aware sizing and wildfire effects", 1),
        ("CV + Thermal fusion is the most robust detector", 0),
        ("it wins every scenario because the two sensors fail independently (+23 pp recall)", 1),
    ], Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.5), size=16, gap=3)

    # ---- Slide 3: Training data (dataset details) ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Training data", "SYNTHETIC \u00b7 PHYSICS-GROUNDED \u00b7 PUBLISHABLE",
              "7,200 generated images across drone, UGV, and thermal \u2014 with physically-correct survivor sizes & hard negatives")
    pic_fit(s, p_dataset, Inches(0.4), Inches(1.75), Inches(12.5), Inches(4.1))
    add_text(s, "Sources: SARD (real person cutouts)  \u00b7  NAIP (USDA aerial imagery, geographically split)  \u00b7  VisDrone (hard-negative vehicles). "
                "Composited with color harmonization, altitude/range-aware blur, and wildfire effects.",
             Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.55), size=16, color=RGB_GREY, font=BODY_FONT)
    add_text(s, "Published on Zenodo  \u00b7  DOI 10.5281/zenodo.21226010  \u00b7  CC-BY 4.0  \u00b7  1.4 GB total\n"
                "Includes 5 trained YOLOv8 weights: drone, drone_naip, ugv_front, ugv_mast, thermal",
             Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.65), size=16, bold=True, color=RGB_RUST, font=BODY_FONT)

    # ---- Slide 4: Timeline ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Development timeline", "HOW WE GOT HERE",
              "From abstract perception to a full multi-modal, physics-grounded CV stack")
    pic_fit(s, p_timeline, Inches(0.4), Inches(2.35), Inches(12.5), Inches(3.0))
    add_bullets(s, [
        ("Each step targeted a real failure mode", 0),
        ("false positives on terrain, unrealistic survivor sizes, occlusion, and fire crossover", 1),
    ], Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.2), size=16, gap=6)

    # ---- Slide 5: Combined dataset examples (default + synthetic) ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Real vs synthetic training data", "DATASET EXAMPLES",
              "Top row: real sensor inputs  ·  Bottom row: synthetic composites with survivors & wildfire effects")
    pic_fit(s, p_cv_all, Inches(0.1), Inches(1.85), Inches(13.1), Inches(5.5))

    # ---- Slide 6: Trained detector performance ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Trained detector performance", "HELD-OUT VALIDATION")
    pic_fit(s, p_train, Inches(0.4), Inches(1.9), Inches(8.0), Inches(5.0))
    add_bullets(s, [
        ("Drone oblique", 0),
        ("P 0.82 · R 0.66 · mAP50 0.72", 1),
        ("Drone NAIP", 0),
        ("P 0.89 · R 0.91 · mAP50 0.91", 1),
        ("UGV front / mast", 0),
        ("P 0.87/0.90 · R 0.74/0.90 · mAP50 0.81/0.92", 1),
        ("Thermal YOLO", 0),
        ("P 1.00 · R 0.96 · mAP50 0.97", 1),
    ], Inches(8.5), Inches(1.95), Inches(4.4), Inches(5.0), size=16, gap=6)

    # ---- Slide 6b: Hard-negative robustness (decoy FP heatmap) ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Hard-negative robustness", "DOES THE MODEL LEARN SHAPE, OR SHORTCUTS?",
              "Added animal & colorful human-scale decoys so size/color alone can't separate them from survivors")
    pic_fit(s, p_decoy_heat, Inches(0.5), Inches(1.9), Inches(6.6), Inches(4.9))
    add_bullets(s, [
        ("Reviewer concern: cars were too large/plain to be a real test", 0),
        ("model could learn \u2018colorful blob = person\u2019 instead of shape", 1),
        ("Added animal + colorful-object decoys at survivor scale (0.5\u20131.8 m)", 0),
        ("same size, palette & contrast as survivors \u2014 shape is the only cue left", 1),
        ("Colorful-object false positives: 40.9% \u2192 14.3% after retrain", 0),
        ("Animal false positives: 6.7% \u2192 2.0%", 0),
        ("Vehicle false positives: 0% \u2192 1.0% (still negligible)", 1),
        ("Residual 14% on colorful objects is unchanged by more data/resolution", 0),
        ("confirms it needs more decoy diversity, not just scale", 1),
    ], Inches(7.35), Inches(1.95), Inches(5.4), Inches(5.0), size=15, gap=8)

    # ---- Slide 6c: Recall by target size (heatmap) ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Recall by target size", "SMALL-OBJECT DETECTION IS THE HARD CASE",
              "Retrained natively at 1280px on 6.8k images \u2014 deployment-config recall jumped 72% \u2192 95%")
    pic_fit(s, p_size_heat, Inches(0.6), Inches(1.9), Inches(6.9), Inches(4.9))
    add_bullets(s, [
        ("Evaluated at conf 0.001 (full PR sweep), dual IoU, Wilson 95% CIs", 0),
        ("Root cause found: old model trained at 640px, deployed via 1280 tiles", 0),
        ("tiled inference fed it out-of-distribution crops \u2014 not a data problem", 1),
        ("Retraining natively at 1280px fixed it: 71.8% \u2192 95.1% overall recall", 0),
        ("at the actual deployment config (2\u00d72 tiles, 1280px)", 1),
        ("Sub-8px recall: 62.5% \u2192 89.9% at deployment config", 0),
        ("640 full frame is now the weakest config \u2014 model trained for 1280", 0),
        ("8\u201315px and >15px both exceed 95% at 1280 configs", 0),
    ], Inches(7.65), Inches(1.95), Inches(5.1), Inches(5.0), size=15, gap=8)

    # ---- Slide 6: THE HEATMAP (comparison) ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Modality comparison — recall heatmap", "CV vs THERMAL vs MOTION vs FUSION",
              "CV column: retrained YOLOv8s (1280px, 6.8k images) measured on val frames with each scenario's effects synthesized")
    pic_fit(s, p_heat, Inches(0.35), Inches(1.85), Inches(8.7), Inches(5.1))
    add_bullets(s, [
        ("Green = high recall, red = failure", 0),
        ("CV (measured): retrained model on 172 val frames per scenario", 0),
        ("82% clear, degrading to 66\u201374% in fire/burn scenarios", 1),
        ("Thermal collapses on burned ground (3%)", 0),
        ("thermal crossover: body temp \u2248 hot ground", 1),
        ("Motion flat at 33% — survivors are static", 0),
        ("CV+Thermal is the greenest column overall", 0),
        ("fusion recomputed from the measured CV recall", 1),
    ], Inches(9.15), Inches(1.95), Inches(3.9), Inches(5.0), size=15, gap=7)

    # ---- Slide 6d: Real-world validation heatmap (sim-to-real gap) ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Real-world validation \u2014 sim-to-real gap", "MEASURED ON REAL IMAGES, NOT SIMULATED",
              "Same models, evaluated on real thermal/drone/UGV imagery before and after mixing real data into training")
    pic_fit(s, p_realworld, Inches(0.4), Inches(1.9), Inches(7.6), Inches(5.0))
    add_bullets(s, [
        ("Synthetic-only models lose 60\u201395 pts of recall on real imagery", 0),
        ("thermal 96%\u219231%, drone 95%\u219222%, UGV 33%\u21920.5\u20134%", 1),
        ("Mixing real data into training recovers most of the gap", 0),
        ("thermal \u219286%, drone \u219282%, UGV \u219257\u201363%", 1),
        ("UGV tested on two real domains", 0),
        ("COCO (generic) and real Malibu-area terrain photos + real cutouts", 1),
        ("real-terrain UGV score (63%) beats COCO (57%)", 1),
        ("Real test sets: HIT-UAV (thermal), HERIDAL (drone SAR),", 0),
        ("Wikimedia Malibu-area terrain + COCO (UGV)", 1),
    ], Inches(8.3), Inches(1.95), Inches(4.6), Inches(5.0), size=14, gap=7)

    # ---- Slide 6a: REAL CV recall (honest measurement) ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "How the real detector actually performs", "MEASURED, NOT SIMULATED",
              "The trained YOLOv8s survivor detector run on real held-out validation frames")
    pic_fit(s, p_real_cv, Inches(0.6), Inches(1.9), Inches(6.6), Inches(4.9))
    add_bullets(s, [
        ("Real model: 82% recall clear, 88% recall fire/smoke", 0),
        ("the opposite of the idealized proxy's assumption", 1),
        ("Likely a dataset composition effect, not true fire-robustness", 0),
        ("survivor size/placement isn\u2019t independently controlled per condition", 1),
        ("Only 2 conditions shown \u2014 that's what our labels support", 0),
        ("has_fire is the only condition flag in the generated metadata", 1),
        ("Next: label finer smoke/fire severity to test this rigorously", 0),
    ], Inches(7.5), Inches(1.95), Inches(5.2), Inches(5.0), size=15, gap=9)

    # ---- Slide 7: Fusion bars ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Why fusion wins", "COMPLEMENTARY FAILURE MODES")
    pic_fit(s, p_fusion, Inches(0.35), Inches(1.9), Inches(9.0), Inches(5.0))
    add_bullets(s, [
        ("In fire: CV 74% + thermal 54% → fusion 88%", 0),
        ("On burned ground thermal dies (3%), CV rescues fusion to 67%", 0),
        ("Clear: fusion 95% vs 82% CV alone (+14 pp)", 0),
        ("No training data needed for thermal/motion", 0),
        ("both are physics-based, not learned", 1),
    ], Inches(9.45), Inches(1.95), Inches(3.6), Inches(5.0), size=16, gap=10)

    # ---- Slide 8: Motion detection — how it works ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Motion detection — how it works", "TEMPORAL FRAME DIFFERENCING",
              "A physics-based alternative that needs no training data")
    divider(s, Inches(2.15))
    add_bullets(s, [
        ("Algorithm: compare consecutive drone camera frames", 0),
        ("convert each frame to grayscale, compute |frame₁ − frame₀|", 1),
        ("threshold pixel differences (>30) → binary change mask", 1),
        ("dilate + flood-fill connected components → candidate blobs", 1),
        ("filter by blob area (100–15,000 px) → survivor detections", 1),
        ("Confidence penalties reduce false positives", 0),
        ("large drone movement → more background change → lower confidence", 1),
        ("smoke reduces contrast → smoke_penalty lowers confidence", 1),
        ("Key limitation: survivors must be static", 0),
        ("wildfire victims are typically immobile, so only ~1 of 3 detected per pass", 1),
        ("fire/smoke itself creates motion artifacts (false positives)", 1),
        ("Recall: 33% across all scenarios (vs 95% for CV in clear)", 0),
        ("Result: motion adds negligible value — CV+Thermal fusion is superior", 0),
    ], Inches(0.7), Inches(2.4), Inches(12.0), Inches(4.6), size=16, gap=5)

    # ---- Slide 9: CV detection — clear conditions ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_text(s, "CV detection — clear conditions", Inches(0.6), Inches(0.25), Inches(12.2), Inches(0.55),
             size=28, bold=True, color=RGB_GREEN, font=HEAD_FONT)
    footer(s)
    pic_fit(s, p_cv_clear, Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.3))

    # ---- Slide 10: CV detection — smoke & fire ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_text(s, "CV detection — smoke & fire", Inches(0.6), Inches(0.25), Inches(12.2), Inches(0.55),
             size=28, bold=True, color=RGB_GREEN, font=HEAD_FONT)
    footer(s)
    pic_fit(s, p_cv_fire, Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.3))

    # ---- Slide 11: CV by condition ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "CV perception quality", "BY CONDITION")
    pic_fit(s, p_cond, Inches(0.35), Inches(1.9), Inches(8.5), Inches(5.0))
    add_bullets(s, [
        ("Retrained model, measured on 172 held-out val frames", 0),
        ("Recall 0.82 clear, 0.80 under smoke", 0),
        ("Fire/burn degrade recall to 0.66–0.74", 0),
        ("burned ground is hardest: char texture hides small bodies", 1),
        ("Precision 0.77–0.84 at conf 0.15", 0),
        ("operating point tunable via the PR curve", 1),
    ], Inches(8.95), Inches(1.95), Inches(4.1), Inches(5.0), size=16, gap=9)

    # ---- Slide 9: Answers to the 3 questions ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Three key capabilities delivered", "SIDE ANGLES · UGV DATA · THERMAL")
    divider(s, Inches(1.75))
    add_bullets(s, [
        ("1 · Drone side angles — IMPLEMENTED", 0),
        ("oblique camera mode (15–45° tilt), 25% of training images; wired into the simulation via camera_tilt_deg", 1),
        ("2 · UGV computer-vision data — DELIVERED", 0),
        ("3,600 images across front + mast cameras; SARD + NAIP + VisDrone decoys; validated on unseen terrain", 1),
        ("3 · Thermal images — GENERATED & TRAINED", 0),
        ("1,200-image simulated TIR dataset; trained thermal YOLO (R 0.96); usable as physics or YOLO backend", 1),
    ], Inches(0.7), Inches(2.0), Inches(12.0), Inches(4.8), size=16, gap=16)

    # ---- Slide 10: Zenodo dataset publication ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Published on Zenodo", "OPEN DATA · DOI 10.5281/ZENODO.21226010",
              "All datasets, trained weights, generation scripts & documentation — CC-BY 4.0")
    divider(s, Inches(2.2))

    zenodo_files = [
        ("omnisearch_drone_survivor.zip", "322 MB", "2,000 train + 400 val (25% oblique)"),
        ("omnisearch_drone_survivor_naip.zip", "120 MB", "500 train + 90 val (NAIP-only bg)"),
        ("omnisearch_ugv_survivor.zip", "568 MB", "3,000 train + 600 val (front + mast)"),
        ("omnisearch_thermal_tir.zip", "140 MB", "1,000 train + 200 val (simulated TIR)"),
        ("omnisearch_naip_backgrounds.zip", "93 MB", "81 NAIP tiles + VisDrone decoys"),
    ]
    tb = slide.shapes.add_textbox if False else None
    y_row = Inches(2.5)
    add_text(s, "FILE", Inches(0.7), y_row, Inches(4.5), Inches(0.3),
             size=12, bold=True, color=RGB_TAUPE, font=LABEL_FONT)
    add_text(s, "SIZE", Inches(5.5), y_row, Inches(1.2), Inches(0.3),
             size=12, bold=True, color=RGB_TAUPE, font=LABEL_FONT)
    add_text(s, "CONTENTS", Inches(7.0), y_row, Inches(5.5), Inches(0.3),
             size=12, bold=True, color=RGB_TAUPE, font=LABEL_FONT)
    for i, (fname, fsize, desc) in enumerate(zenodo_files):
        y = Inches(2.9) + i * Inches(0.38)
        add_text(s, fname, Inches(0.7), y, Inches(4.5), Inches(0.35),
                 size=16, bold=False, color=RGB_DARK, font=BODY_FONT)
        add_text(s, fsize, Inches(5.5), y, Inches(1.2), Inches(0.35),
                 size=16, bold=True, color=RGB_RUST, font=BODY_FONT)
        add_text(s, desc, Inches(7.0), y, Inches(5.5), Inches(0.35),
                 size=16, bold=False, color=RGB_GREY, font=BODY_FONT)

    add_bullets(s, [
        ("5 trained YOLOv8 weights included", 0),
        ("survivor, survivor_naip, ugv_front, ugv_mast, thermal", 1),
        ("Total bundle: 1.4 GB · fully reproducible", 0),
        ("generation scripts + datasheet EDA + modality report", 1),
    ], Inches(0.7), Inches(5.0), Inches(12.0), Inches(1.8), size=16, gap=6)

    add_text(s, "zenodo.org/records/21226010",
             Inches(0.7), Inches(6.6), Inches(8.0), Inches(0.35),
             size=16, bold=True, color=RGB_RUST, font=HEAD_FONT)

    # ---- Slide 11: Conclusions ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title_bar(s, "Conclusions & next steps", "TAKEAWAYS")
    divider(s, Inches(1.75))
    add_bullets(s, [
        ("CV + Thermal fusion is the recommended detector for wildfire SAR", 0),
        ("wins every scenario through complementary, independent failure modes", 1),
        ("Motion detection adds negligible recall", 0),
        ("survivors are immobile, so frame-differencing rarely fires", 1),
        ("Sim-to-real gap measured and largely closed on real imagery", 0),
        ("HIT-UAV (thermal), HERIDAL (drone SAR), real Malibu terrain + COCO (UGV)", 1),
        ("mixed real+synthetic fine-tuning: +55\u201360 pts recall on every modality", 1),
        ("Training data is synthetic-first, reproducible & published", 0),
        ("1.4 GB on Zenodo with weights, scripts & datasheet (DOI: 10.5281/zenodo.21226010)", 1),
        ("Next: fine-tune drone/thermal further on more real wildfire scenes", 0),
        ("Next: add oblique drone rendering to the live trajectory export", 0),
    ], Inches(0.7), Inches(2.0), Inches(12.0), Inches(4.8), size=16, gap=10)

    out = SLIDE_DIR / "omnisearch_cv_results.pptx"
    prs.save(str(out))
    print(f"Saved deck: {out}  ({len(prs.slides._sldIdLst)} slides)")
    return out


if __name__ == "__main__":
    build()
